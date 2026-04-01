from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color palette
DARK_BG = RGBColor(0x1B, 0x1B, 0x2F)
ACCENT_BLUE = RGBColor(0x00, 0x7A, 0xCC)
ACCENT_CYAN = RGBColor(0x00, 0xBC, 0xD4)
ACCENT_GREEN = RGBColor(0x4C, 0xAF, 0x50)
ACCENT_ORANGE = RGBColor(0xFF, 0x98, 0x00)
ACCENT_PURPLE = RGBColor(0x9C, 0x27, 0xB0)
ACCENT_RED = RGBColor(0xF4, 0x43, 0x36)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
SECTION_BG = RGBColor(0x0D, 0x47, 0xA1)
CARD_BG = RGBColor(0x2A, 0x2A, 0x45)


def add_dark_background(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG


def add_accent_bar(slide, color=ACCENT_BLUE, height=Inches(0.06)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_slide_number(slide, number, total):
    txBox = slide.shapes.add_textbox(
        Inches(12.2), Inches(7.0), Inches(1.0), Inches(0.4)
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"{number}/{total}"
    p.font.size = Pt(10)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.RIGHT


def add_section_badge(slide, text, color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(0.4), Inches(2.0), Inches(0.35)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(11)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER


def add_title_text(
    slide,
    text,
    left,
    top,
    width,
    height,
    size=36,
    color=WHITE,
    bold=True,
    align=PP_ALIGN.LEFT,
):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return tf


def add_bullet_list(
    slide,
    items,
    left,
    top,
    width,
    height,
    size=18,
    color=WHITE,
    spacing=Pt(8),
    bullet_color=ACCENT_CYAN,
):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = spacing
        p.level = 0
    return tf


def add_card(
    slide,
    left,
    top,
    width,
    height,
    title,
    items,
    title_color=ACCENT_CYAN,
    bg_color=CARD_BG,
):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    shape.shadow.inherit = False

    # Title
    txBox = slide.shapes.add_textbox(
        left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), Inches(0.4)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.color.rgb = title_color
    p.font.bold = True

    # Items
    txBox2 = slide.shapes.add_textbox(
        left + Inches(0.2),
        top + Inches(0.55),
        width - Inches(0.4),
        height - Inches(0.7),
    )
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(13)
        p.font.color.rgb = LIGHT_GRAY
        p.space_after = Pt(4)


def add_speaker_notes(slide, text):
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


TOTAL_SLIDES = 20

# ============================================================
# SLIDE 1: Title Slide
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_dark_background(slide)
add_accent_bar(slide, ACCENT_BLUE, Inches(0.08))

# Top accent line
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(2), Inches(2.5), Inches(9.333), Inches(0.04)
)
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_CYAN
shape.line.fill.background()

add_title_text(
    slide,
    "Module 1",
    Inches(2),
    Inches(2.7),
    Inches(9.333),
    Inches(0.8),
    size=24,
    color=ACCENT_CYAN,
    bold=False,
    align=PP_ALIGN.CENTER,
)
add_title_text(
    slide,
    "Generative AI & Prompt Engineering",
    Inches(2),
    Inches(3.3),
    Inches(9.333),
    Inches(1.0),
    size=44,
    color=WHITE,
    bold=True,
    align=PP_ALIGN.CENTER,
)
add_title_text(
    slide,
    "Foundations for the Applied AI Era",
    Inches(2),
    Inches(4.2),
    Inches(9.333),
    Inches(0.6),
    size=20,
    color=LIGHT_GRAY,
    bold=False,
    align=PP_ALIGN.CENTER,
)

# Bottom accent line
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(2), Inches(5.2), Inches(9.333), Inches(0.04)
)
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_CYAN
shape.line.fill.background()

add_title_text(
    slide,
    "Applied AI Course  |  Part 1: Foundations",
    Inches(2),
    Inches(5.5),
    Inches(9.333),
    Inches(0.5),
    size=14,
    color=LIGHT_GRAY,
    bold=False,
    align=PP_ALIGN.CENTER,
)

add_speaker_notes(
    slide,
    (
        "Welcome to Module 1 of the Applied AI course. This module covers the foundational concepts of Generative AI, "
        "including what it is, how LLMs work, prompt engineering techniques, data analysis with LLMs, and the ethical "
        "considerations surrounding AI. By the end of this module, students will understand the capabilities and limitations "
        "of generative AI and be able to craft effective prompts for various tasks."
    ),
)
add_slide_number(slide, 1, TOTAL_SLIDES)

# ============================================================
# SLIDE 2: Agenda / Module Overview
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_background(slide)
add_accent_bar(slide, ACCENT_BLUE)
add_section_badge(slide, "OVERVIEW", ACCENT_BLUE)
add_title_text(
    slide, "Module Agenda", Inches(0.5), Inches(1.0), Inches(12), Inches(0.8), size=36
)

sections = [
    (
        "1.1",
        "What is Generative AI?",
        "LLM architecture, capabilities, and limitations",
        ACCENT_BLUE,
    ),
    (
        "1.2",
        "Prompt Engineering",
        "Techniques: zero-shot, few-shot, CoT, ReAct, ToT",
        ACCENT_CYAN,
    ),
    (
        "1.3",
        "Data Analysis with LLMs",
        "SQL generation, visualization, statistical analysis",
        ACCENT_GREEN,
    ),
    (
        "1.4",
        "Ethics & the Future of AI",
        "Bias, fairness, regulations, emerging trends",
        ACCENT_PURPLE,
    ),
    (
        "",
        "Enterprise Scenario",
        "Real-world production application of GenAI",
        ACCENT_ORANGE,
    ),
]

for i, (num, title, desc, color) in enumerate(sections):
    y = Inches(2.0) + Inches(i * 1.0)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(0.7), Inches(0.7)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = num if num else "E"
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(8)

    add_title_text(
        slide,
        title,
        Inches(1.8),
        y + Inches(0.0),
        Inches(5),
        Inches(0.4),
        size=20,
        color=WHITE,
        bold=True,
    )
    add_title_text(
        slide,
        desc,
        Inches(1.8),
        y + Inches(0.35),
        Inches(8),
        Inches(0.35),
        size=14,
        color=LIGHT_GRAY,
        bold=False,
    )

add_speaker_notes(
    slide,
    (
        "This module is structured into four main sections plus an enterprise scenario. We start with the fundamentals of "
        "generative AI and LLM architecture, then move into practical prompt engineering techniques. Section 3 covers how "
        "LLMs can assist with data analysis tasks. Section 4 addresses the critical ethical dimensions. Finally, we tie "
        "everything together with a production-level enterprise scenario demonstrating how these concepts apply in real business."
    ),
)
add_slide_number(slide, 2, TOTAL_SLIDES)

# ============================================================
# SLIDE 3: What is Generative AI?
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_background(slide)
add_accent_bar(slide, ACCENT_BLUE)
add_section_badge(slide, "SECTION 1.1", ACCENT_BLUE)
add_title_text(
    slide,
    "What is Generative AI?",
    Inches(0.5),
    Inches(1.0),
    Inches(12),
    Inches(0.8),
    size=36,
)

add_title_text(
    slide,
    (
        "AI systems that create novel content — text, images, code, audio, video — "
        "by learning patterns from training data."
    ),
    Inches(0.5),
    Inches(1.9),
    Inches(12),
    Inches(0.7),
    size=18,
    color=LIGHT_GRAY,
    bold=False,
)

# Two cards side by side
add_card(
    slide,
    Inches(0.5),
    Inches(2.9),
    Inches(5.8),
    Inches(3.8),
    "Discriminative AI",
    [
        "Learns decision boundaries between classes",
        "Classification and prediction tasks",
        "Examples: spam detection, image recognition",
        "Models: SVM, Logistic Regression, BERT",
    ],
    title_color=ACCENT_ORANGE,
)

add_card(
    slide,
    Inches(7.0),
    Inches(2.9),
    Inches(5.8),
    Inches(3.8),
    "Generative AI",
    [
        "Learns underlying data distribution",
        "Creates new content that resembles training data",
        "Examples: text generation, image synthesis",
        "Models: GPT, DALL-E, Stable Diffusion, Claude",
    ],
    title_color=ACCENT_CYAN,
)

add_speaker_notes(
    slide,
    (
        "Generative AI represents a paradigm shift from traditional discriminative AI. While discriminative models learn "
        "boundaries between classes (e.g., spam vs. not spam), generative models learn the full data distribution and can "
        "sample from it to produce entirely new outputs. This enables applications like conversational AI, code generation, "
        "creative writing, and multimodal content creation. The key distinction: discriminative = 'What is this?' vs. generative = 'Create something new.'"
    ),
)
add_slide_number(slide, 3, TOTAL_SLIDES)

# ============================================================
# SLIDE 4: LLM Architecture
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_background(slide)
add_accent_bar(slide, ACCENT_BLUE)
add_section_badge(slide, "SECTION 1.1", ACCENT_BLUE)
add_title_text(
    slide,
    "Large Language Model Architecture",
    Inches(0.5),
    Inches(1.0),
    Inches(12),
    Inches(0.8),
    size=36,
)

components = [
    ("Tokenizer", "Converts text to subword tokens (BPE, WordPiece)", ACCENT_BLUE),
    ("Embedding Layer", "Maps tokens to dense vector representations", ACCENT_CYAN),
    (
        "Transformer Blocks",
        "Self-attention + feed-forward networks (stacked)",
        ACCENT_GREEN,
    ),
    ("Positional Encoding", "Injects sequence order information", ACCENT_ORANGE),
    (
        "Attention Mechanism",
        "Weighted relevance between all token pairs",
        ACCENT_PURPLE,
    ),
    ("Output Head", "Probability distribution over vocabulary", ACCENT_RED),
]

for i, (name, desc, color) in enumerate(components):
    row = i // 3
    col = i % 3
    x = Inches(0.5) + Inches(col * 4.2)
    y = Inches(2.1) + Inches(row * 2.4)

    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.9), Inches(2.0)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = color
    shape.line.width = Pt(2)

    add_title_text(
        slide,
        name,
        x + Inches(0.2),
        y + Inches(0.2),
        Inches(3.5),
        Inches(0.4),
        size=18,
        color=color,
        bold=True,
    )
    add_title_text(
        slide,
        desc,
        x + Inches(0.2),
        y + Inches(0.7),
        Inches(3.5),
        Inches(1.1),
        size=14,
        color=LIGHT_GRAY,
        bold=False,
    )

add_speaker_notes(
    slide,
    (
        "LLMs are built on the transformer architecture introduced in 'Attention Is All You Need' (Vaswani et al., 2017). "
        "The key innovation is the self-attention mechanism which computes relationships between all token pairs simultaneously, "
        "enabling parallel processing and capturing long-range dependencies. The architecture consists of: tokenization (converting "
        "text to subword units), embedding (mapping to dense vectors), stacked transformer blocks (each with multi-head attention "
        "and feed-forward layers with residual connections and layer norm), positional encoding (injecting order since attention "
        "is permutation-invariant), and an output head that produces a probability distribution over the vocabulary for next-token prediction."
    ),
)
add_slide_number(slide, 4, TOTAL_SLIDES)

# ============================================================
# SLIDE 5: LLM Capabilities
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_background(slide)
add_accent_bar(slide, ACCENT_GREEN)
add_section_badge(slide, "SECTION 1.1", ACCENT_BLUE)
add_title_text(
    slide,
    "LLM Capabilities",
    Inches(0.5),
    Inches(1.0),
    Inches(12),
    Inches(0.8),
    size=36,
)

capabilities = [
    ("Text Generation", "Articles, stories, code, emails, marketing copy"),
    ("Summarization", "Extractive and abstractive summarization of documents"),
    ("Translation", "Cross-lingual text translation with context awareness"),
    ("Question Answering", "Open-domain and closed-domain QA systems"),
    ("Code Generation", "Function synthesis, bug fixing, code explanation"),
    ("Reasoning", "Mathematical, logical, and commonsense reasoning"),
    ("Instruction Following", "Task completion from natural language instructions"),
]

for i, (cap, desc) in enumerate(capabilities):
    row = i // 4
    col = i % 4
    x = Inches(0.3) + Inches(col * 3.2)
    y = Inches(2.1) + Inches(row * 2.6)
    w = Inches(3.0)

    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, Inches(2.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.fill.background()

    # Icon-like number
    add_title_text(
        slide,
        str(i + 1),
        x + Inches(0.2),
        y + Inches(0.15),
        Inches(0.4),
        Inches(0.4),
        size=20,
        color=ACCENT_GREEN,
        bold=True,
    )
    add_title_text(
        slide,
        cap,
        x + Inches(0.2),
        y + Inches(0.55),
        Inches(2.6),
        Inches(0.4),
        size=16,
        color=WHITE,
        bold=True,
    )
    add_title_text(
        slide,
        desc,
        x + Inches(0.2),
        y + Inches(1.0),
        Inches(2.6),
        Inches(1.0),
        size=12,
        color=LIGHT_GRAY,
        bold=False,
    )

add_speaker_notes(
    slide,
    (
        "Modern LLMs have a remarkably broad capability surface. They can generate coherent long-form text, summarize "
        "documents, translate between languages, answer questions, write and debug code, perform multi-step reasoning, "
        "and follow complex instructions. These capabilities make them versatile tools for automation, augmentation, and "
        "creative applications across virtually every industry. The key is understanding which capability to leverage for "
        "a given task and how to prompt the model effectively to achieve the desired output."
    ),
)
add_slide_number(slide, 5, TOTAL_SLIDES)

# ============================================================
# SLIDE 6: LLM Limitations
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_background(slide)
add_accent_bar(slide, ACCENT_RED)
add_section_badge(slide, "SECTION 1.1", ACCENT_BLUE)
add_title_text(
    slide,
    "LLM Limitations & Mitigations",
    Inches(0.5),
    Inches(1.0),
    Inches(12),
    Inches(0.8),
    size=36,
)

limitations = [
    (
        "Hallucinations",
        "Plausible but incorrect information",
        "RAG, verification, grounding",
        ACCENT_RED,
    ),
    (
        "Knowledge Cutoff",
        "Training data has a fixed date",
        "Web search, fine-tuning, RAG",
        ACCENT_ORANGE,
    ),
    (
        "Context Window",
        "Limited input/output token capacity",
        "Chunking, summarization, sliding windows",
        ACCENT_BLUE,
    ),
    (
        "Bias",
        "Reflects biases in training data",
        "Filtering, diverse data, auditing",
        ACCENT_PURPLE,
    ),
    (
        "Reasoning Errors",
        "Struggles with complex multi-step logic",
        "Chain-of-thought, tool use, decomposition",
        ACCENT_CYAN,
    ),
    (
        "Non-Determinism",
        "Non-deterministic outputs (temperature)",
        "Temperature=0, seed values, structured output",
        ACCENT_GREEN,
    ),
    (
        "Cost & Latency",
        "API calls can be expensive and slow",
        "Caching, smaller models, batching",
        ACCENT_ORANGE,
    ),
]

headers = ["Limitation", "Description", "Mitigation"]
col_widths = [Inches(2.5), Inches(4.5), Inches(5.0)]
col_starts = [Inches(0.5), Inches(3.0), Inches(7.5)]
header_y = Inches(2.0)

for j, (header, cw, cx) in enumerate(zip(headers, col_widths, col_starts)):
    add_title_text(
        slide,
        header,
        cx,
        header_y,
        cw,
        Inches(0.4),
        size=16,
        color=ACCENT_CYAN,
        bold=True,
    )

for i, (lim, desc, mit, color) in enumerate(limitations):
    y = Inches(2.5) + Inches(i * 0.65)
    # Subtle row indicator
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.3), y + Inches(0.05), Inches(0.08), Inches(0.35)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

    add_title_text(
        slide,
        lim,
        col_starts[0],
        y,
        col_widths[0],
        Inches(0.35),
        size=14,
        color=WHITE,
        bold=True,
    )
    add_title_text(
        slide,
        desc,
        col_starts[1],
        y,
        col_widths[1],
        Inches(0.35),
        size=13,
        color=LIGHT_GRAY,
        bold=False,
    )
    add_title_text(
        slide,
        mit,
        col_starts[2],
        y,
        col_widths[2],
        Inches(0.35),
        size=13,
        color=ACCENT_GREEN,
        bold=False,
    )

add_speaker_notes(
    slide,
    (
        "Understanding LLM limitations is critical for production deployment. Hallucinations are the most dangerous — "
        "the model generates confident-sounding but incorrect information. Knowledge cutoff means the model doesn't know "
        "about events after its training date. Context windows limit how much information can be processed at once. Bias "
        "reflects societal prejudices in training data. Reasoning errors occur with complex multi-step logic. Non-determinism "
        "means the same prompt can yield different outputs. Cost and latency are real operational concerns at scale. "
        "Each limitation has specific mitigation strategies that should be implemented in production systems."
    ),
)
add_slide_number(slide, 6, TOTAL_SLIDES)

# ============================================================
# SLIDE 7: Section Divider - Prompt Engineering
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_background(slide)
bg_shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height
)
bg_shape.fill.solid()
bg_shape.fill.fore_color.rgb = SECTION_BG
bg_shape.line.fill.background()

add_title_text(
    slide,
    "Section 1.2",
    Inches(0),
    Inches(2.5),
    prs.slide_width,
    Inches(0.6),
    size=20,
    color=ACCENT_CYAN,
    bold=False,
    align=PP_ALIGN.CENTER,
)
add_title_text(
    slide,
    "Prompt Engineering",
    Inches(0),
    Inches(3.1),
    prs.slide_width,
    Inches(1.0),
    size=48,
    color=WHITE,
    bold=True,
    align=PP_ALIGN.CENTER,
)
add_title_text(
    slide,
    "The art and science of guiding LLMs toward desired outputs",
    Inches(0),
    Inches(4.2),
    prs.slide_width,
    Inches(0.6),
    size=18,
    color=LIGHT_GRAY,
    bold=False,
    align=PP_ALIGN.CENTER,
)

add_speaker_notes(
    slide,
    (
        "Prompt engineering is arguably the most important practical skill for working with LLMs. It involves designing, "
        "testing, and refining the inputs we give to language models to reliably produce the outputs we need. This section "
        "covers foundational through advanced prompting techniques that every AI practitioner should master."
    ),
)
add_slide_number(slide, 7, TOTAL_SLIDES)

# ============================================================
# SLIDE 8: Prompt Engineering Core Principles
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_background(slide)
add_accent_bar(slide, ACCENT_CYAN)
add_section_badge(slide, "SECTION 1.2", ACCENT_CYAN)
add_title_text(
    slide,
    "Core Principles of Prompt Engineering",
    Inches(0.5),
    Inches(1.0),
    Inches(12),
    Inches(0.8),
    size=36,
)

principles = [
    (
        "Be Specific",
        "Clear, unambiguous instructions produce better results",
        ACCENT_BLUE,
    ),
    (
        "Provide Context",
        "Background information improves relevance and accuracy",
        ACCENT_CYAN,
    ),
    (
        "Use Examples",
        "Demonstrations guide the model's output format and style",
        ACCENT_GREEN,
    ),
    (
        "Structure Prompts",
        "Use delimiters, sections, and formatting for clarity",
        ACCENT_ORANGE,
    ),
    (
        "Iterate & Test",
        "Test and refine prompts systematically with A/B testing",
        ACCENT_PURPLE,
    ),
]

for i, (title, desc, color) in enumerate(principles):
    y = Inches(2.1) + Inches(i * 1.0)
    # Number badge
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(0.8), y + Inches(0.05), Inches(0.55), Inches(0.55)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = str(i + 1)
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    add_title_text(
        slide,
        title,
        Inches(1.6),
        y,
        Inches(3),
        Inches(0.4),
        size=20,
        color=WHITE,
        bold=True,
    )
    add_title_text(
        slide,
        desc,
        Inches(1.6),
        y + Inches(0.35),
        Inches(10),
        Inches(0.35),
        size=14,
        color=LIGHT_GRAY,
        bold=False,
    )

add_speaker_notes(
    slide,
    (
        "These five principles form the foundation of effective prompt engineering. Being specific eliminates ambiguity. "
        "Providing context gives the model the background it needs. Using examples (few-shot) dramatically improves output "
        "quality and format adherence. Structuring prompts with delimiters and sections helps the model distinguish between "
        "instructions and data. Finally, iteration is key — treat prompt development like software development with version "
        "control, testing, and continuous improvement."
    ),
)
add_slide_number(slide, 8, TOTAL_SLIDES)

# ============================================================
# SLIDE 9: Prompting Techniques
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_background(slide)
add_accent_bar(slide, ACCENT_CYAN)
add_section_badge(slide, "SECTION 1.2", ACCENT_CYAN)
add_title_text(
    slide,
    "Prompting Techniques",
    Inches(0.5),
    Inches(1.0),
    Inches(12),
    Inches(0.8),
    size=36,
)

techniques = [
    (
        "Zero-Shot",
        "No examples provided. Model relies on pre-trained knowledge.\n\nClassify sentiment: 'The service was terrible.'\nSentiment:",
        ACCENT_BLUE,
    ),
    (
        "Few-Shot",
        "2-5 examples demonstrate expected format.\n\n'Great!' → Positive\n'Terrible' → Negative\n'Okay' → Neutral\n'Late delivery' →",
        ACCENT_CYAN,
    ),
    (
        "Chain-of-Thought",
        "Step-by-step reasoning before final answer.\n\n'Sell 1/4 of 120, then 1/3 of remaining.\nStep 1: 120 × 1/4 = 30 sold...\nStep 2: 90 × 1/3 = 30 sold...\nAnswer: 60'",
        ACCENT_GREEN,
    ),
    (
        "Role-Based",
        "Assign persona to guide behavior.\n\n'You are a senior security engineer.\nReview this code for vulnerabilities...'",
        ACCENT_ORANGE,
    ),
]

for i, (name, example, color) in enumerate(techniques):
    x = Inches(0.3) + Inches(i * 3.25)
    y = Inches(2.1)
    w = Inches(3.05)

    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, Inches(4.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = color
    shape.line.width = Pt(2)

    # Header bar
    hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Inches(0.5))
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = color
    hdr.line.fill.background()
    tf = hdr.text_frame
    p = tf.paragraphs[0]
    p.text = name
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    add_title_text(
        slide,
        example,
        x + Inches(0.15),
        y + Inches(0.65),
        w - Inches(0.3),
        Inches(4.0),
        size=12,
        color=LIGHT_GRAY,
        bold=False,
    )

add_speaker_notes(
    slide,
    (
        "These four foundational techniques represent increasing levels of guidance. Zero-shot is simplest but least reliable "
        "for complex tasks. Few-shot dramatically improves output quality by providing input-output examples. Chain-of-thought "
        "is essential for reasoning tasks — it forces the model to show its work, which improves accuracy on math, logic, and "
        "multi-step problems. Role-based prompting leverages the model's ability to adopt personas, which is useful for "
        "specialized tasks like code review, medical analysis, or legal document drafting. In practice, these techniques are "
        "often combined for optimal results."
    ),
)
add_slide_number(slide, 9, TOTAL_SLIDES)

# ============================================================
# SLIDE 10: Advanced Prompt Patterns
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_background(slide)
add_accent_bar(slide, ACCENT_PURPLE)
add_section_badge(slide, "SECTION 1.2", ACCENT_CYAN)
add_title_text(
    slide,
    "Advanced Prompt Patterns",
    Inches(0.5),
    Inches(1.0),
    Inches(12),
    Inches(0.8),
    size=36,
)

add_card(
    slide,
    Inches(0.3),
    Inches(2.0),
    Inches(4.0),
    Inches(4.8),
    "ReAct (Reasoning + Acting)",
    [
        "Alternates between Thought, Action, Observation",
        "Integrates tool use (search, calculator, DB)",
        "Handles tasks requiring external knowledge",
        "Foundation for agentic AI systems",
        "Example: Thought → search() → observe → answer",
    ],
    title_color=ACCENT_BLUE,
)

add_card(
    slide,
    Inches(4.6),
    Inches(2.0),
    Inches(4.0),
    Inches(4.8),
    "Tree of Thoughts (ToT)",
    [
        "Explores multiple reasoning paths simultaneously",
        "Evaluates and selects the best path",
        "Uses BFS/DFS/beam search strategies",
        "Ideal for complex planning and design tasks",
        "Example: 3 approaches → evaluate pros/cons → pick best",
    ],
    title_color=ACCENT_GREEN,
)

add_card(
    slide,
    Inches(8.9),
    Inches(2.0),
    Inches(4.0),
    Inches(4.8),
    "Self-Consistency & Chaining",
    [
        "Self-Consistency: generate N answers, majority vote",
        "Improves reliability for reasoning tasks",
        "Prompt Chaining: sequential multi-step prompts",
        "Each step's output feeds into the next",
        "Example: extract → summarize → analyze",
    ],
    title_color=ACCENT_ORANGE,
)

add_speaker_notes(
    slide,
    (
        "Advanced prompt patterns go beyond single-turn prompting. ReAct (Yao et al., 2022) is the foundation for modern "
        "agentic AI — it combines reasoning traces with tool actions, enabling LLMs to interact with external systems. "
        "Tree of Thoughts (Yao et al., 2023) extends chain-of-thought by exploring multiple reasoning branches and selecting "
        "the best one, similar to how humans brainstorm solutions. Self-consistency (Wang et al., 2022) generates multiple "
        "reasoning paths and uses majority voting to improve accuracy. Prompt chaining decomposes complex tasks into "
        "sequential steps, where each step's output becomes the next step's input — this is essentially building pipelines "
        "with LLMs."
    ),
)
add_slide_number(slide, 10, TOTAL_SLIDES)

# ============================================================
# SLIDE 11: Prompt Best Practices
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_background(slide)
add_accent_bar(slide, ACCENT_CYAN)
add_section_badge(slide, "SECTION 1.2", ACCENT_CYAN)
add_title_text(
    slide,
    "Prompt Best Practices",
    Inches(0.5),
    Inches(1.0),
    Inches(12),
    Inches(0.8),
    size=36,
)

practices = [
    ("Use Delimiters", 'Separate instructions from data using """, ---, ###'),
    ("Specify Format", "Define expected output structure (JSON, XML, markdown)"),
    ("Set Constraints", "Limit output scope: 'In 3 bullet points...'"),
    ("Provide Context", "Give background: 'Analyzing Q3 financial data...'"),
    ("Use Examples", "Show input-output pairs for few-shot learning"),
    ("Iterative Refinement", "A/B test prompt variants systematically"),
    ("Handle Edge Cases", "Specify behavior for uncertain/missing inputs"),
    ("Structured Output", "Use response_format for programmatic integration"),
]

for i, (title, desc) in enumerate(practices):
    row = i // 4
    col = i % 4
    x = Inches(0.3) + Inches(col * 3.25)
    y = Inches(2.1) + Inches(row * 2.6)

    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.05), Inches(2.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.fill.background()

    add_title_text(
        slide,
        title,
        x + Inches(0.2),
        y + Inches(0.2),
        Inches(2.65),
        Inches(0.4),
        size=16,
        color=ACCENT_CYAN,
        bold=True,
    )
    add_title_text(
        slide,
        desc,
        x + Inches(0.2),
        y + Inches(0.7),
        Inches(2.65),
        Inches(1.2),
        size=13,
        color=LIGHT_GRAY,
        bold=False,
    )

add_speaker_notes(
    slide,
    (
        "These best practices should be applied systematically. Delimiters are critical for preventing prompt injection and "
        "helping the model distinguish between instructions and data. Specifying output format ensures parseable results. "
        "Setting constraints prevents overly verbose responses. Providing context improves relevance. Examples are the single "
        "most impactful technique for improving output quality. Iterative refinement with A/B testing mirrors software testing "
        "practices. Handling edge cases prevents unexpected failures. Structured output with response_format guarantees "
        "machine-parseable JSON, which is essential for production pipelines."
    ),
)
add_slide_number(slide, 11, TOTAL_SLIDES)

# ============================================================
# SLIDE 12: Section Divider - Data Analysis
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg_shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height
)
bg_shape.fill.solid()
bg_shape.fill.fore_color.rgb = RGBColor(0x1B, 0x5E, 0x20)
bg_shape.line.fill.background()

add_title_text(
    slide,
    "Section 1.3",
    Inches(0),
    Inches(2.5),
    prs.slide_width,
    Inches(0.6),
    size=20,
    color=ACCENT_GREEN,
    bold=False,
    align=PP_ALIGN.CENTER,
)
add_title_text(
    slide,
    "Data Analysis with LLMs",
    Inches(0),
    Inches(3.1),
    prs.slide_width,
    Inches(1.0),
    size=48,
    color=WHITE,
    bold=True,
    align=PP_ALIGN.CENTER,
)
add_title_text(
    slide,
    "Transforming raw data into actionable insights with AI assistance",
    Inches(0),
    Inches(4.2),
    prs.slide_width,
    Inches(0.6),
    size=18,
    color=LIGHT_GRAY,
    bold=False,
    align=PP_ALIGN.CENTER,
)

add_speaker_notes(
    slide,
    (
        "LLMs are powerful assistants for data analysis tasks. They can generate cleaning code, write SQL queries from "
        "natural language, create visualization code, and even suggest statistical approaches. This section covers practical "
        "use cases where LLMs accelerate the data analysis workflow."
    ),
)
add_slide_number(slide, 12, TOTAL_SLIDES)

# ============================================================
# SLIDE 13: Data Analysis Use Cases
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_background(slide)
add_accent_bar(slide, ACCENT_GREEN)
add_section_badge(slide, "SECTION 1.3", ACCENT_GREEN)
add_title_text(
    slide,
    "Data Analysis Use Cases with LLMs",
    Inches(0.5),
    Inches(1.0),
    Inches(12),
    Inches(0.8),
    size=36,
)

use_cases = [
    (
        "Data Cleaning",
        "Identify missing values, outliers, inconsistencies. Generate pandas cleaning code.",
        ACCENT_BLUE,
    ),
    (
        "SQL Generation",
        "Convert natural language questions into SQL queries. Schema-aware query building.",
        ACCENT_CYAN,
    ),
    (
        "Visualization",
        "Generate matplotlib/seaborn/plotly code for charts and dashboards.",
        ACCENT_GREEN,
    ),
    (
        "Statistical Analysis",
        "Suggest tests, compute correlations, recommend modeling approaches.",
        ACCENT_ORANGE,
    ),
    (
        "Report Writing",
        "Create narrative summaries of analysis findings and trends.",
        ACCENT_PURPLE,
    ),
    (
        "Data Transformation",
        "Generate code for reshaping, aggregating, merging datasets.",
        ACCENT_RED,
    ),
]

for i, (title, desc, color) in enumerate(use_cases):
    row = i // 3
    col = i % 3
    x = Inches(0.5) + Inches(col * 4.2)
    y = Inches(2.1) + Inches(row * 2.5)

    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.9), Inches(2.1)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = color
    shape.line.width = Pt(2)

    # Color bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.08), Inches(2.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

    add_title_text(
        slide,
        title,
        x + Inches(0.25),
        y + Inches(0.2),
        Inches(3.4),
        Inches(0.4),
        size=18,
        color=color,
        bold=True,
    )
    add_title_text(
        slide,
        desc,
        x + Inches(0.25),
        y + Inches(0.7),
        Inches(3.4),
        Inches(1.2),
        size=13,
        color=LIGHT_GRAY,
        bold=False,
    )

add_speaker_notes(
    slide,
    (
        "These six use cases represent the highest-value applications of LLMs in data analysis. Data cleaning is often "
        "80% of a data scientist's work — LLMs can identify issues and generate cleaning code. SQL generation democratizes "
        "data access, allowing non-technical users to query databases. Visualization code generation accelerates dashboard "
        "creation. Statistical analysis assistance helps choose appropriate tests and interpret results. Report writing "
        "automates the narrative around analysis. Data transformation generates the reshape/aggregate code that's tedious "
        "to write manually. Critical caveat: always validate LLM-generated code against actual data before production use."
    ),
)
add_slide_number(slide, 13, TOTAL_SLIDES)

# ============================================================
# SLIDE 14: Section Divider - Ethics
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg_shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height
)
bg_shape.fill.solid()
bg_shape.fill.fore_color.rgb = RGBColor(0x4A, 0x14, 0x8C)
bg_shape.line.fill.background()

add_title_text(
    slide,
    "Section 1.4",
    Inches(0),
    Inches(2.5),
    prs.slide_width,
    Inches(0.6),
    size=20,
    color=ACCENT_PURPLE,
    bold=False,
    align=PP_ALIGN.CENTER,
)
add_title_text(
    slide,
    "Ethics, Considerations & the Future",
    Inches(0),
    Inches(3.1),
    prs.slide_width,
    Inches(1.0),
    size=48,
    color=WHITE,
    bold=True,
    align=PP_ALIGN.CENTER,
)
add_title_text(
    slide,
    "Building responsible AI systems for a better future",
    Inches(0),
    Inches(4.2),
    prs.slide_width,
    Inches(0.6),
    size=18,
    color=LIGHT_GRAY,
    bold=False,
    align=PP_ALIGN.CENTER,
)

add_speaker_notes(
    slide,
    (
        "As AI systems become more powerful and pervasive, ethical considerations are not optional — they are essential. "
        "This section covers bias, transparency, privacy, safety, the regulatory landscape, and future trends that every "
        "AI practitioner must understand."
    ),
)
add_slide_number(slide, 14, TOTAL_SLIDES)

# ============================================================
# SLIDE 15: Bias and Responsible AI
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_background(slide)
add_accent_bar(slide, ACCENT_PURPLE)
add_section_badge(slide, "SECTION 1.4", ACCENT_PURPLE)
add_title_text(
    slide,
    "Bias, Fairness & Responsible AI",
    Inches(0.5),
    Inches(1.0),
    Inches(12),
    Inches(0.8),
    size=36,
)

# Bias types
bias_types = [
    ("Representation", "Under/over-representation of groups"),
    ("Historical", "Past societal inequalities reflected"),
    ("Measurement", "Flawed feature selection / proxies"),
    ("Aggregation", "One-size-fits-all models"),
    ("Evaluation", "Biased benchmark datasets"),
]

for i, (name, desc) in enumerate(bias_types):
    x = Inches(0.5) + Inches(i * 2.55)
    y = Inches(2.0)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.35), Inches(1.8)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = ACCENT_RED
    shape.line.width = Pt(1)
    add_title_text(
        slide,
        name,
        x + Inches(0.15),
        y + Inches(0.15),
        Inches(2.05),
        Inches(0.35),
        size=14,
        color=ACCENT_RED,
        bold=True,
    )
    add_title_text(
        slide,
        desc,
        x + Inches(0.15),
        y + Inches(0.6),
        Inches(2.05),
        Inches(1.0),
        size=12,
        color=LIGHT_GRAY,
        bold=False,
    )

# Responsible AI pillars
pillars = ["Fairness", "Transparency", "Accountability", "Privacy", "Safety"]
pillar_colors = [ACCENT_BLUE, ACCENT_CYAN, ACCENT_GREEN, ACCENT_ORANGE, ACCENT_PURPLE]

add_title_text(
    slide,
    "Responsible AI Framework",
    Inches(0.5),
    Inches(4.1),
    Inches(12),
    Inches(0.5),
    size=20,
    color=WHITE,
    bold=True,
)

for i, (pillar, color) in enumerate(zip(pillars, pillar_colors)):
    x = Inches(0.5) + Inches(i * 2.55)
    y = Inches(4.7)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.35), Inches(0.55)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = pillar
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

add_title_text(
    slide,
    (
        "Bias detection & mitigation  ·  Model documentation & explainability  ·  Human oversight  ·  "
        "Data minimization & consent  ·  Content filtering & monitoring"
    ),
    Inches(0.5),
    Inches(5.5),
    Inches(12),
    Inches(1.2),
    size=13,
    color=LIGHT_GRAY,
    bold=False,
)

add_speaker_notes(
    slide,
    (
        "Bias in AI is a systemic issue. Representation bias occurs when certain groups are underrepresented in training data. "
        "Historical bias perpetuates past inequalities. Measurement bias uses flawed proxies (e.g., zip code for creditworthiness). "
        "Aggregation bias applies one model to diverse populations. Evaluation bias uses skewed benchmarks. "
        "The Responsible AI Framework addresses these through five pillars: Fairness (bias detection, diverse data, audits), "
        "Transparency (model documentation, explainability), Accountability (human oversight, error reporting), "
        "Privacy (data minimization, consent management), and Safety (content filtering, monitoring, incident response)."
    ),
)
add_slide_number(slide, 15, TOTAL_SLIDES)

# ============================================================
# SLIDE 16: Regulatory Landscape & Future
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_background(slide)
add_accent_bar(slide, ACCENT_PURPLE)
add_section_badge(slide, "SECTION 1.4", ACCENT_PURPLE)
add_title_text(
    slide,
    "Regulatory Landscape & Future Trends",
    Inches(0.5),
    Inches(1.0),
    Inches(12),
    Inches(0.8),
    size=36,
)

# Regulations
regs = [
    ("EU AI Act", "EU", "Risk-based classification, transparency"),
    ("NIST AI RMF", "US", "Voluntary risk management framework"),
    ("China AI Regs", "China", "Algorithm registration, content controls"),
    ("UK AI Safety", "UK", "Pro-innovation, sector-specific"),
    ("Canada AIDA", "Canada", "Risk management, accountability"),
]

add_title_text(
    slide,
    "Key Regulations",
    Inches(0.5),
    Inches(1.9),
    Inches(5),
    Inches(0.4),
    size=18,
    color=WHITE,
    bold=True,
)
for i, (name, region, desc) in enumerate(regs):
    y = Inches(2.4) + Inches(i * 0.55)
    add_title_text(
        slide,
        f"{name} ({region})",
        Inches(0.5),
        y,
        Inches(3),
        Inches(0.35),
        size=13,
        color=ACCENT_ORANGE,
        bold=True,
    )
    add_title_text(
        slide,
        desc,
        Inches(3.5),
        y,
        Inches(3),
        Inches(0.35),
        size=13,
        color=LIGHT_GRAY,
        bold=False,
    )

# Future trends
trends = [
    ("Multimodal Models", "Unified text, image, audio, video"),
    ("Agentic AI", "Autonomous planning & execution"),
    ("Small Language Models", "Efficient, domain-specific, edge"),
    ("Neuro-Symbolic AI", "Neural + symbolic reasoning"),
    ("AI Alignment", "Systems following human intent"),
]

add_title_text(
    slide,
    "Emerging Trends",
    Inches(7.0),
    Inches(1.9),
    Inches(5),
    Inches(0.4),
    size=18,
    color=WHITE,
    bold=True,
)
for i, (name, desc) in enumerate(trends):
    y = Inches(2.4) + Inches(i * 0.55)
    add_title_text(
        slide,
        name,
        Inches(7.0),
        y,
        Inches(3),
        Inches(0.35),
        size=13,
        color=ACCENT_CYAN,
        bold=True,
    )
    add_title_text(
        slide,
        desc,
        Inches(10.0),
        y,
        Inches(3),
        Inches(0.35),
        size=13,
        color=LIGHT_GRAY,
        bold=False,
    )

add_speaker_notes(
    slide,
    (
        "The regulatory landscape is evolving rapidly. The EU AI Act is the most comprehensive, using a risk-based classification "
        "system (unacceptable, high, limited, minimal risk). NIST AI RMF provides voluntary guidance for US organizations. "
        "China requires algorithm registration. The UK takes a pro-innovation approach. Canada's AIDA focuses on accountability. "
        "Future trends include multimodal models (GPT-4V, Gemini), agentic AI that can plan and execute autonomously, "
        "small language models for efficiency and privacy (Phi, Gemma), neuro-symbolic AI combining neural and logical reasoning, "
        "and AI alignment research to ensure systems follow human intent."
    ),
)
add_slide_number(slide, 16, TOTAL_SLIDES)

# ============================================================
# SLIDE 17: Enterprise Scenario - Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg_shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height
)
bg_shape.fill.solid()
bg_shape.fill.fore_color.rgb = RGBColor(0xE6, 0x51, 0x00)
bg_shape.line.fill.background()

add_title_text(
    slide,
    "Enterprise Scenario",
    Inches(0),
    Inches(2.2),
    prs.slide_width,
    Inches(0.6),
    size=20,
    color=WHITE,
    bold=False,
    align=PP_ALIGN.CENTER,
)
add_title_text(
    slide,
    "AI-Powered Customer Support",
    Inches(0),
    Inches(2.8),
    prs.slide_width,
    Inches(1.0),
    size=48,
    color=WHITE,
    bold=True,
    align=PP_ALIGN.CENTER,
)
add_title_text(
    slide,
    "Enterprise Insurance Corp — Automating 70% of Tier-1 Support Tickets",
    Inches(0),
    Inches(4.0),
    prs.slide_width,
    Inches(0.6),
    size=18,
    color=RGBColor(0xFF, 0xCC, 0x80),
    bold=False,
    align=PP_ALIGN.CENTER,
)

add_speaker_notes(
    slide,
    (
        "Let's bring all these concepts together with a real enterprise scenario. We'll examine how a large insurance company "
        "deploys generative AI to transform their customer support operations — covering prompt engineering, data analysis, "
        "ethical guardrails, and production considerations."
    ),
)
add_slide_number(slide, 17, TOTAL_SLIDES)

# ============================================================
# SLIDE 18: Enterprise Scenario - Details
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_background(slide)
add_accent_bar(slide, ACCENT_ORANGE)
add_section_badge(slide, "ENTERPRISE SCENARIO", ACCENT_ORANGE)
add_title_text(
    slide,
    "Enterprise Insurance Corp — Production Architecture",
    Inches(0.5),
    Inches(1.0),
    Inches(12),
    Inches(0.8),
    size=34,
)

# Scenario cards
add_card(
    slide,
    Inches(0.3),
    Inches(2.0),
    Inches(4.0),
    Inches(2.2),
    "Business Problem",
    [
        "50,000+ support tickets/month",
        "Average resolution: 4.2 hours",
        "$2.8M annual support cost",
        "Customer satisfaction: 62%",
    ],
    title_color=ACCENT_RED,
)

add_card(
    slide,
    Inches(4.6),
    Inches(2.0),
    Inches(4.0),
    Inches(2.2),
    "GenAI Solution",
    [
        "LLM-powered ticket classification",
        "Automated response generation (RAG)",
        "Intelligent routing to specialists",
        "Real-time sentiment analysis",
    ],
    title_color=ACCENT_CYAN,
)

add_card(
    slide,
    Inches(8.9),
    Inches(2.0),
    Inches(4.0),
    Inches(2.2),
    "Techniques Applied",
    [
        "Few-shot prompting for classification",
        "RAG for policy-specific answers",
        "Chain-of-thought for claim assessment",
        "Structured output for CRM integration",
    ],
    title_color=ACCENT_GREEN,
)

# Results
add_title_text(
    slide,
    "Production Results",
    Inches(0.5),
    Inches(4.5),
    Inches(12),
    Inches(0.4),
    size=20,
    color=WHITE,
    bold=True,
)

results = [
    ("70%", "Tickets\nAutomated", ACCENT_BLUE),
    ("1.2 hrs", "Avg Resolution\n(down from 4.2)", ACCENT_CYAN),
    ("$1.8M", "Annual Savings\n(64% cost reduction)", ACCENT_GREEN),
    ("89%", "Customer\nSatisfaction", ACCENT_ORANGE),
    ("99.7%", "System\nUptime", ACCENT_PURPLE),
]

for i, (value, label, color) in enumerate(results):
    x = Inches(0.3) + Inches(i * 2.6)
    y = Inches(5.0)

    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.4), Inches(1.8)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = color
    shape.line.width = Pt(2)

    add_title_text(
        slide,
        value,
        x + Inches(0.1),
        y + Inches(0.15),
        Inches(2.2),
        Inches(0.5),
        size=28,
        color=color,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_title_text(
        slide,
        label,
        x + Inches(0.1),
        y + Inches(0.8),
        Inches(2.2),
        Inches(0.8),
        size=12,
        color=LIGHT_GRAY,
        bold=False,
        align=PP_ALIGN.CENTER,
    )

add_speaker_notes(
    slide,
    (
        "Enterprise Insurance Corp processes 50,000+ support tickets monthly with an average resolution time of 4.2 hours "
        "at $2.8M annual cost and only 62% customer satisfaction. Their GenAI solution uses: (1) Few-shot prompting to "
        "classify tickets into categories (claims, billing, policy questions, complaints), (2) RAG to ground responses in "
        "actual policy documents reducing hallucinations, (3) Chain-of-thought for complex claim assessments, and "
        "(4) Structured output (JSON) to integrate with their Salesforce CRM. Results after 6 months: 70% automation rate, "
        "resolution time down to 1.2 hours, $1.8M savings, 89% satisfaction, 99.7% uptime. Key lesson: human-in-the-loop "
        "for high-stakes decisions (claim denials) and continuous prompt optimization based on feedback data."
    ),
)
add_slide_number(slide, 18, TOTAL_SLIDES)

# ============================================================
# SLIDE 19: Key Takeaways
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_background(slide)
add_accent_bar(slide, ACCENT_BLUE)
add_section_badge(slide, "SUMMARY", ACCENT_BLUE)
add_title_text(
    slide, "Key Takeaways", Inches(0.5), Inches(1.0), Inches(12), Inches(0.8), size=36
)

takeaways = [
    (
        "Generative AI creates novel content using learned data distributions; LLMs are transformer-based models with broad capabilities and important limitations",
        ACCENT_BLUE,
    ),
    (
        "Prompt engineering is essential — master zero-shot, few-shot, chain-of-thought, ReAct, and structured output techniques",
        ACCENT_CYAN,
    ),
    (
        "LLMs accelerate data analysis through SQL generation, visualization code, and statistical assistance — but always validate outputs",
        ACCENT_GREEN,
    ),
    (
        "Ethical AI requires addressing bias, ensuring transparency, protecting privacy, and complying with evolving regulations",
        ACCENT_PURPLE,
    ),
    (
        "Production deployment demands reliability, monitoring, cost management, security, and human oversight for high-stakes decisions",
        ACCENT_ORANGE,
    ),
]

for i, (text, color) in enumerate(takeaways):
    y = Inches(2.0) + Inches(i * 1.05)

    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), y, Inches(0.5), Inches(0.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = str(i + 1)
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    add_title_text(
        slide,
        text,
        Inches(1.3),
        y + Inches(0.05),
        Inches(11.5),
        Inches(0.5),
        size=16,
        color=WHITE,
        bold=False,
    )

add_speaker_notes(
    slide,
    (
        "These five takeaways encapsulate the core learning objectives of Module 1. (1) Understanding generative AI fundamentals "
        "— what it can and cannot do. (2) Prompt engineering as the primary interface skill for LLMs. (3) Practical data analysis "
        "applications with the critical caveat of validation. (4) The non-negotiable importance of ethical AI practices. "
        "(5) The operational realities of production deployment. Students should feel confident in these foundations before "
        "proceeding to LangChain, RAG, and agentic systems in subsequent modules."
    ),
)
add_slide_number(slide, 19, TOTAL_SLIDES)

# ============================================================
# SLIDE 20: Thank You / Q&A
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_background(slide)

shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(2), Inches(2.5), Inches(9.333), Inches(0.04)
)
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_CYAN
shape.line.fill.background()

add_title_text(
    slide,
    "Thank You",
    Inches(0),
    Inches(2.8),
    prs.slide_width,
    Inches(1.0),
    size=48,
    color=WHITE,
    bold=True,
    align=PP_ALIGN.CENTER,
)
add_title_text(
    slide,
    "Questions & Discussion",
    Inches(0),
    Inches(3.8),
    prs.slide_width,
    Inches(0.6),
    size=24,
    color=ACCENT_CYAN,
    bold=False,
    align=PP_ALIGN.CENTER,
)

shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(2), Inches(4.6), Inches(9.333), Inches(0.04)
)
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_CYAN
shape.line.fill.background()

add_title_text(
    slide,
    "Next: Module 2 — LangChain & Building Blocks",
    Inches(0),
    Inches(5.0),
    prs.slide_width,
    Inches(0.5),
    size=16,
    color=LIGHT_GRAY,
    bold=False,
    align=PP_ALIGN.CENTER,
)
add_title_text(
    slide,
    "Applied AI Course  |  Part 1: Foundations",
    Inches(0),
    Inches(5.5),
    prs.slide_width,
    Inches(0.5),
    size=14,
    color=LIGHT_GRAY,
    bold=False,
    align=PP_ALIGN.CENTER,
)

add_speaker_notes(
    slide,
    (
        "This concludes Module 1. Encourage students to review the concepts, interview questions, and quiz materials. "
        "The interview questions cover beginner to advanced levels and are excellent preparation for technical assessments. "
        "The quiz with 20 questions and detailed explanations serves as a self-assessment tool. Next module covers LangChain "
        "and building blocks for constructing LLM-powered applications."
    ),
)
add_slide_number(slide, 20, TOTAL_SLIDES)

# Save
output_path = r"D:\Jay Rathod\Tutorials\Applied AI\course-content\part-1-foundations\module-1-generative-ai\Module-1-Generative-AI-and-Prompt-Engineering.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
